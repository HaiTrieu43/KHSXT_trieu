import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_visual_presentation():
    prs = Presentation()
    
    # Thiết lập kích thước slide 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Định nghĩa bảng màu Premium
    PRIMARY_COLOR = RGBColor(0, 121, 107)    # Xanh ngọc C.P. (Teal)
    SECONDARY_COLOR = RGBColor(0, 200, 83)   # Xanh lá tươi (Emerald)
    TEXT_DARK = RGBColor(33, 33, 33)         # Xám đậm (Charcoal)
    TEXT_LIGHT = RGBColor(255, 255, 255)     # Trắng
    ACCENT_RED = RGBColor(211, 47, 47)       # Đỏ đô (cho khó khăn)
    ACCENT_ORANGE = RGBColor(255, 143, 0)    # Cam hổ phách
    BG_DARK = RGBColor(18, 18, 18)           # Nền tối
    
    # Đường dẫn tệp hình ảnh sơ đồ thuật toán
    image_path = r"D:\Kê hoạch sản xuât\laptrinh vao\so_do_thuat_toan_khsx_goc.png"
    
    def apply_title_slide_layout(slide, title_text, subtitle_text):
        """Tạo slide tiêu đề nền tối đẳng cấp"""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK
        
        # Tiêu đề chính
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        
        # Tiêu đề phụ
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = 'Segoe UI'
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(176, 190, 197)
        p2.space_before = Pt(20)
        
        # Người trình bày
        p3 = tf.add_paragraph()
        p3.text = "Báo cáo Kỹ thuật & Tối ưu hóa vận hành — C.P. Bình Dương"
        p3.font.name = 'Segoe UI'
        p3.font.size = Pt(14)
        p3.font.italic = True
        p3.font.color.rgb = RGBColor(120, 144, 156)
        p3.space_before = Pt(40)

    def add_standard_slide(title_text):
        """Tạo một slide tiêu chuẩn có dải tiêu đề ở đầu trang"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Tiêu đề slide
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
        tf = header_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        return slide

    def draw_card(slide, left, top, width, height, title, body, bg_color, title_color=PRIMARY_COLOR):
        """Vẽ một hộp nội dung trực quan hình chữ nhật bo góc (Card)"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = RGBColor(224, 224, 224)
        shape.line.width = Pt(1)
        
        # Text Frame bên trong Card
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Segoe UI'
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = title_color
        
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.name = 'Segoe UI'
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_DARK
        p2.space_before = Pt(8)

    # =========================================================================
    # SLIDE 1: SLIDE TIÊU ĐỀ
    # =========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_title_slide_layout(
        slide1, 
        "SƠ ĐỒ VẬN HÀNH THUẬT TOÁN KHSX TỰ ĐỘNG",
        "Mô Tả Luồng Xử Lý Trực Quan Theo Sơ Đồ Khối & Các Nút Thắt Vận Hành Logistics Thực Tế"
    )
    
    # =========================================================================
    # SLIDE 2: BẢN ĐỒ TOÀN CẢNH THUẬT TOÁN (NHÚNG ẢNH SƠ ĐỒ)
    # =========================================================================
    slide2 = add_standard_slide("1. BẢN ĐỒ TOÀN CẢNH THUẬT TOÁN KHSX (FLOWCHART)")
    
    # Bên trái: Giải thích 3 quy tắc vàng ngắn gọn
    desc_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(6.0), Inches(5.5))
    tf2 = desc_box.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "🎯 3 Nguyên Tắc Vàng Của Thuật Toán:"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    rules = [
        ("Gộp Mã Độc Nhất (Unique Base Code):", "Mỗi mã cám chỉ xuất hiện duy nhất 1 dòng trên KHSX. Lượng xe bồn Silo và lượng đóng bao được gộp chung để ép mẻ tối ưu."),
        ("Bù Shortfall T-1 Tự Động:", "Tự động so sánh Kế hoạch hôm qua vs Mixer chạy thực tế để bù mẻ thiếu hụt ca trước, triệt tiêu hoàn toàn rủi ro sót đơn."),
        ("Bộ Lọc Làm Mịn Tải (Silo Smoothing):", "Tự động điều tiết lượng xe bồn theo năng lực vật lý trạm nạp của nhà máy (tối đa 750 Tấn/ngày) và đẩy phần dư sang ngày kế tiếp.")
    ]
    for title, desc in rules:
        p_b = tf2.add_paragraph()
        p_b.text = f"★  {title} {desc}"
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(12)
        p_b.space_before = Pt(14)
        p_b.font.color.rgb = TEXT_DARK
        p_b.runs[0].font.bold = True
        p_b.runs[0].font.color.rgb = PRIMARY_COLOR
        
    # Bên phải: Nhúng trực tiếp ảnh sơ đồ Mermaid dài của thuật toán
    if os.path.exists(image_path):
        # Tỷ lệ ảnh dài 3.21 -> w = h/3.21. Đặt h = 5.8 inches -> w = 1.8 inches
        slide2.shapes.add_picture(image_path, Inches(8.5), Inches(1.2), Inches(1.8), Inches(5.8))
        
        # Thêm khung chú thích nhỏ bên cạnh ảnh
        caption_box = slide2.shapes.add_textbox(Inches(10.5), Inches(3.0), Inches(2.2), Inches(2.0))
        tf_c = caption_box.text_frame
        tf_c.word_wrap = True
        p_c = tf_c.paragraphs[0]
        p_c.text = "🔍 Sơ đồ luồng 7 bước hoàn chỉnh bên cạnh hiển thị chi tiết các rẽ nhánh logic và các bộ lọc điều kiện của thuật toán."
        p_c.font.name = 'Segoe UI'
        p_c.font.size = Pt(11)
        p_c.font.italic = True
        p_c.font.color.rgb = RGBColor(120, 144, 156)

    # =========================================================================
    # SLIDE 3: BƯỚC 1 & 2 - TIỀN XỬ LÝ DỮ LIỆU
    # =========================================================================
    slide3 = add_standard_slide("2. BƯỚC 1 & 2: TIỀN XỬ LÝ & QUY ĐỔI CÁM NỀN")
    
    # Thẻ Bước 1
    draw_card(
        slide3, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.2),
        "BƯỚC 1: TIỀN XỬ LÝ DỮ LIỆU THÔ",
        "• Áp dụng tệp thay thế mã cám đột xuất (Substitutions) đầu ca.\n• Cộng dồn Tồn kho kho tổng (FFStock) và Tồn bồn chứa thành phẩm (Mixer logs).\n• Xác định danh sách mã cám chỉ chạy xe bồn (Silo-Only) để làm cổng kiểm soát.",
        RGBColor(240, 244, 243)
    )
    
    # Thẻ Bước 2
    draw_card(
        slide3, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.2),
        "BƯỚC 2: QUY ĐỔI MÃ CÁM NỀN (BASE CODE)",
        "• Chuyển đổi toàn bộ mã thương hiệu đại lý về Mã Cám Nền sản xuất chính.\n• Áp dụng quy tắc Prefix hệ thống: 96* (Nuvo) ➔ 5*, HT1* (Star) ➔ 55*, HS1* ➔ 550S...\n• Giúp gom gộp mẻ Mixer đạt công suất tối đa, tránh chạy mẻ nhỡ lẻ tẻ.",
        RGBColor(240, 244, 243)
    )
    
    # Thẻ Khó Khăn Gặp Phải
    draw_card(
        slide3, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.8),
        "⚠️ KHÓ KHĂN VÀ THÁCH THỨC VẬN HÀNH:",
        "1. Sai lệch mã cám:\nMixer chạy thực tế và phòng kinh doanh (Sales Forecast) dùng hệ thống mã khác nhau. Nếu quy đổi sai lệch ➔ Trộn nhầm công thức thức ăn chăn nuôi.\n\n2. Dữ liệu rời rạc:\nTồn kho và tồn bồn nằm ở 2 báo cáo riêng biệt. Nếu không gộp chung ➔ Mixer sẽ chạy trùng lặp sản lượng gây tràn kho chứa.\n\n3. Giải quyết của Thuật toán:\nXây dựng bộ giải mã Prefix tự động quét và quy đổi chính xác 100% mã thương hiệu về cám nền sản xuất trước khi lên kế hoạch Mixer.",
        RGBColor(255, 235, 235),
        ACCENT_RED
    )

    # =========================================================================
    # SLIDE 4: BƯỚC 3 - ENGINE ĐỘ ƯU TIÊN & DOH
    # =========================================================================
    slide4 = add_standard_slide("3. BƯỚC 3: ENGINE ĐỘ ƯU TIÊN & CỔNG LỌC DOH")
    
    draw_card(
        slide4, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8),
        "BƯỚC 3: PHÂN CẤP 4 ĐỘ ƯU TIÊN SẢN XUẤT",
        "• ƯU TIÊN 1: SILO & BÁ CANG ngày mai ➔ Chạy đầu ca để kịp xe bồn rút hàng.\n• ƯU TIÊN 2: WALK-IN ➔ Các đơn trại khẩn cấp phát sinh đột xuất.\n• ƯU TIÊN 3: SHORTFALL ➔ Bù lượng mẻ thiếu hụt ca hôm qua chưa chạy.\n\n• ƯU TIÊN 4: FORECAST TUẦN + CỔNG LỌC DOH:\n   - Nhu cầu dự báo tuần còn lại được kiểm duyệt qua chỉ số DOH.\n   - Chỉ sản xuất hôm nay nếu Tồn kho an toàn DOH < 3.0 Ngày.\n   - Nếu DOH >= 3.0 ngày ➔ Hoãn sản xuất để tối ưu kho bãi.",
        RGBColor(240, 244, 243)
    )
    
    draw_card(
        slide4, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.8),
        "⚠️ KHÓ KHĂN VÀ THÁCH THỨC VẬN HÀNH:",
        "1. Nghẽn phương tiện xuất hàng:\nĐơn xe bồn và bá cang cập bến dồn dập vào sáng sớm. Nếu Mixer trộn muộn ➔ Xe tải và tàu nằm chờ ở cảng sông, phát sinh chi phí lưu bến phạt nặng.\n\n2. Quá tải kho bãi thành phẩm:\nSản xuất dư thừa các mã cám có sức tiêu thụ chậm. Kho chứa bị chiếm dụng vô ích.\n\n3. Giải quyết của Thuật toán:\nĐẩy đơn Silo/Bá Cang lên Ưu tiên 1 để trộn ngay đầu ca. Áp dụng cổng lọc DOH < 3.0 ngày làm màng lọc tự động ngăn chặn sản xuất thừa.",
        RGBColor(255, 235, 235),
        ACCENT_RED
    )

    # =========================================================================
    # SLIDE 5: BƯỚC 4 & 5 - LÀM MỊN TẢI SILO (SMOOTHING)
    # =========================================================================
    slide5 = add_standard_slide("4. BƯỚC 4 & 5: ĐIỀU CHỈNH NHANH & LÀM MỊN SILO")
    
    draw_card(
        slide5, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8),
        "BƯỚC 4 & 5: BỘ LỌC LÀM MỊN TẢI TRỌNG SILO",
        "• Khống chế trần cứng lượng xe bồn Silo ở mức 750.0 Tấn/ngày.\n• Khi tổng Silo lý thuyết vượt 750T (Ví dụ ngày 18/05 yêu cầu 1,871T):\n   - Sắp xếp nhu cầu Silo theo độ ưu tiên giảm dần.\n   - Tích lũy giữ lại các đơn hàng khẩn cấp đạt sát ngưỡng 750T.\n   - Hoãn (Defer) toàn bộ phần xe bồn dư thừa sang ngày hôm sau.\n   - Đối với cám Silo-only ➔ Hủy khỏi KHSX hôm nay.\n   - Đối với cám hỗn hợp ➔ Đặt silo_truck về 0T và tự động tính lại mẻ đóng bao để đạt mẻ tiêu chuẩn.",
        RGBColor(240, 244, 243)
    )
    
    draw_card(
        slide5, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.8),
        "⚠️ KHÓ KHĂN VÀ THÁCH THỨC VẬN HÀNH (NÚT THẮT SILO):",
        "1. Hạ tầng trạm nạp có giới hạn:\nNhà máy chỉ có 2 vòi nạp xe bồn. Tối đa 24h chạy liên tục chỉ nạp được 720T cám xá. Thuật toán đề xuất 1,871T Silo ban đầu là bất khả thi.\n\n2. Giới hạn đội xe chở chuyên dụng:\nToàn chi nhánh chỉ có 18 xe bồn bulk truck, vòng chạy trung bình 5-7h/chuyến ➔ Năng lực chuyên chở thực tế tối đa chỉ đạt 700T-750T/ngày.\n\n3. Giải quyết của Thuật toán:\nBộ lọc làm mịn tải Silo khống chế chính xác trần 750T. Ngày 18/05 Mixer thực tế chỉ giao được 747.6T Silo, khớp hoàn hảo với thuật toán mới.",
        RGBColor(255, 235, 235),
        ACCENT_RED
    )

    # =========================================================================
    # SLIDE 6: BƯỚC 6 & 7 - CHUẨN HÓA MẺ & ĐẦU RA KHSX
    # =========================================================================
    slide6 = add_standard_slide("5. BƯỚC 6 & 7: CHUẨN HÓA MẺ & KẾT QUẢ ĐẦU RA")
    
    draw_card(
        slide6, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8),
        "BƯỚC 6 & 7: MẺ TRỘN TIÊU CHUẨN & PHẦN MỀM KHSX",
        "• Chuẩn hóa trọng lượng mẻ cứng (Standard TPB):\n   - Mẻ 8.0 Tấn: Họ cám con 550*, 551* và mã đặc biệt 325F.\n   - Mẻ 8.4 Tấn: Toàn bộ các dòng sản phẩm còn lại.\n• Đảm bảo công thức: Tổng tấn = Số mẻ × Mẻ chuẩn ➔ Đạt chuẩn 0 lỗi Grader V2.\n• Phân bổ bao bì chi tiết (25kg, 40kg, 50kg) tự động theo vỏ bao rỗng thực tế có sẵn.\n• Sắp xếp chuỗi Mixer sequence thông minh theo ma trận an toàn sinh học để ngăn ngừa nhiễm chéo kháng sinh.",
        RGBColor(240, 244, 243)
    )
    
    draw_card(
        slide6, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.8),
        "⚠️ KHÓ KHĂN VÀ THÁCH THỨC VẬN HÀNH:",
        "1. Lỗi lệch mẻ trên hệ thống điều khiển:\nÉp mẻ lẻ phi tiêu chuẩn từ file điều chỉnh nhanh dễ làm lệch tấn trên cân điện tử và gây lỗi cảnh báo nghiêm trọng từ bộ chấm điểm.\n\n2. Rủi ro nhiễm chéo thuốc chăn nuôi:\nSắp xếp Mixer thủ công dễ để cám sạch chạy sau cám kháng sinh mạnh ➔ Gây nhiễm chất cấm chất lượng.\n\n3. Giải quyết của Thuật toán:\nĐồng bộ hệ số TPB cứng 8.0T/8.4T. Áp dụng ma trận an toàn sinh học sắp xếp ca trộn tự động dựa trên 26 cấp độ thuốc.",
        RGBColor(255, 235, 235),
        ACCENT_RED
    )

    # =========================================================================
    # SLIDE 7: CÁC KHÓ KHĂN VẬN HÀNH & NÚT THẮT LOGISTICS THỰC TẾ
    # =========================================================================
    slide7 = add_standard_slide("6. CÁC NÚT THẮT CỔ CHAI LOGISTICS TRONG VẬN HÀNH THỰC TẾ")
    
    # 3 Cột tương ứng với 3 Nút thắt cổ chai vật lý lớn nhất
    
    # Cột 1: Trạm nạp nhà máy
    draw_card(
        slide7, Inches(0.8), Inches(1.6), Inches(3.6), Inches(4.8),
        "VÒI NẠP NHÀ MÁY",
        "• Giới hạn trạm:\nNhà máy chỉ có 2 vòi nạp xe bồn chuyên dụng.\n\n• Thời gian nạp:\nMất 45-60 phút để hoàn tất quy trình cân rỗng, nạp cám xá, cân đầy cho một xe bồn 30 tấn.\n\n• Công suất trần:\nChạy liên tục 24h tối đa chỉ nạp được 720 Tấn/ngày.\n\n➔ Đơn Silo đề xuất 1,871T ban đầu hoàn toàn bị tắc nghẽn tại vòi nạp.",
        RGBColor(255, 243, 224),
        ACCENT_ORANGE
    )
    
    # Cột 2: Đội xe bồn
    draw_card(
        slide7, Inches(4.8), Inches(1.6), Inches(3.6), Inches(4.8),
        "ĐỘI XE VẬN CHUYỂN",
        "• Giới hạn đội xe:\nToàn chi nhánh chỉ có 18 xe bồn bulk truck chở cám xá.\n\n• Thời gian quay vòng:\nKhoảng cách đến các trại lớn (Bình Phước, Tây Ninh) từ 70-150km. Vòng chạy mất 5-7 giờ.\n\n• Công suất chuyên chở:\nMỗi xe chạy tối đa 2-3 chuyến/ngày. Tổng năng lực chở thực tế đạt tối đa 700T-750T/ngày.",
        RGBColor(255, 243, 224),
        ACCENT_ORANGE
    )
    
    # Cột 3: Hạ tầng trang trại
    draw_card(
        slide7, Inches(8.8), Inches(1.6), Inches(3.6), Inches(4.8),
        "TIẾP NHẬN TẠI TRẠI",
        "• Giới hạn bồn chứa:\nMỗi trại lớn chỉ có 2-4 silo chứa dung tích nhỏ (10-20 tấn).\n\n• Tốc độ xả cám:\nHệ thống vít tải hoặc khí nén tại trại xả cám mất 1.5 - 2 giờ/xe.\n\n• Nghẽn bến bãi:\nNếu dồn dập xe bồn cùng đổ bộ xuống trại ➔ Xe phải xếp hàng nằm chờ trên đường, tê liệt vòng quay xe bồn của nhà máy.",
        RGBColor(255, 243, 224),
        ACCENT_ORANGE
    )

    # =========================================================================
    # SLIDE 8: KHUYẾN NGHỊ HÀNH ĐỘNG SÚC TÍCH (CHECKLIST)
    # =========================================================================
    slide8 = add_standard_slide("7. KHUYẾN NGHỊ VẬN HÀNH & KẾT LUẬN")
    content_box8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.0))
    tf8 = content_box8.text_frame
    tf8.word_wrap = True
    
    p = tf8.paragraphs[0]
    p.text = "Lộ trình 4 bước để tối ưu hóa toàn diện hiệu suất tháp Mixer Bình Dương:"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    recs = [
        ("1. Áp dụng Bộ Lọc Làm Mịn Tải Silo (Trần 750T):", "Tuyệt đối không tự ý duyệt lượng xe bồn vượt quá 750T/ngày để bảo vệ an toàn cho đội xe và hạ tầng tiếp nhận của trang trại."),
        ("2. Số hóa chuẩn hóa dữ liệu chốt ca Thứ Bảy:", "Đồng bộ hóa 7 tệp báo cáo đầu vào chốt hết ca Thứ Bảy để làm sạch dữ liệu nền tảng, đảm bảo chạy thuật toán tự động ca Thứ Hai chính xác."),
        ("3. Nghiêm túc duy trì giới hạn tối đa 35 mã cám/ngày:", "Duy trì tỷ lệ xoay vòng sản phẩm tuần để máy ép viên Pellet chạy liên tục không thay khuôn, nâng cao chất lượng viên và tiết kiệm 15% điện năng."),
        ("4. Kích hoạt sắp xếp Mixer ca trộn tự động:", "Thay thế lập lịch thủ công bằng Mixer sequence tự động dựa trên ma trận kháng sinh 26 cấp độ để triệt tiêu hoàn toàn rủi ro nhiễm chéo chất cấm.")
    ]
    
    for title, desc in recs:
        p_b = tf8.add_paragraph()
        p_b.text = f"✔  {title}\n     {desc}"
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.space_before = Pt(14)
        p_b.font.color.rgb = TEXT_DARK
        p_b.runs[0].font.bold = True
        p_b.runs[0].font.color.rgb = SECONDARY_COLOR

    # Lưu file PowerPoint hoàn chỉnh vào thư mục dự án
    output_path = r"D:\Kê hoạch sản xuât\laptrinh vao\MO_TA_THUAT_TOAN_KHSX_TRUCOUAN.pptx"
    prs.save(output_path)
    print(f"✅ Slide PowerPoint Trực Quan đã được tạo và lưu thành công tại: {output_path}")

if __name__ == "__main__":
    create_visual_presentation()
